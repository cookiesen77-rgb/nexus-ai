"""
专业 PPTX 构建器 - 基于 banana-slides 的设计

功能:
- 精确字体计算（使用 PIL ImageFont）
- CJK 字体支持（中日韩文字）
- 多颜色文本支持
- 自动换行和字体适配
- HTML 表格解析
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from io import BytesIO
import base64

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# 尝试导入 PIL
try:
    from PIL import Image, ImageFont, ImageDraw
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    logger.warning("PIL not available, using fallback font size calculation")


class PPTXBuilder:
    """
    专业 PPTX 构建器
    
    基于 banana-slides 的设计，提供：
    - 精确的字体大小计算
    - CJK 字符支持
    - 自动换行
    - 专业的布局
    """
    
    # 标准幻灯片尺寸 (16:9)
    DEFAULT_SLIDE_WIDTH_INCHES = 13.333  # 16:9 宽度
    DEFAULT_SLIDE_HEIGHT_INCHES = 7.5    # 16:9 高度
    
    # DPI 转换
    DEFAULT_DPI = 96
    
    # python-pptx 尺寸限制
    MAX_SLIDE_WIDTH_INCHES = 56.0
    MAX_SLIDE_HEIGHT_INCHES = 56.0
    MIN_SLIDE_WIDTH_INCHES = 1.0
    MIN_SLIDE_HEIGHT_INCHES = 1.0
    
    # 全局字体大小限制
    MIN_FONT_SIZE = 6
    MAX_FONT_SIZE = 200
    
    # 默认字体大小
    DEFAULT_TITLE_FONT_SIZE = 44
    DEFAULT_BODY_FONT_SIZE = 24
    DEFAULT_SUBTITLE_FONT_SIZE = 18
    
    # 默认内边距
    DEFAULT_MARGIN_INCHES = 0.5
    
    # 字体缓存
    _font_cache: Dict[float, Any] = {}
    
    # 尝试查找系统字体
    SYSTEM_FONTS = [
        "/System/Library/Fonts/PingFang.ttc",  # macOS
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",  # Linux
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "C:\\Windows\\Fonts\\msyh.ttc",  # Windows
        "C:\\Windows\\Fonts\\simsun.ttc",
    ]
    
    _system_font_path: Optional[str] = None
    
    @classmethod
    def _find_system_font(cls) -> Optional[str]:
        """查找可用的系统字体"""
        if cls._system_font_path is not None:
            return cls._system_font_path
        
        for font_path in cls.SYSTEM_FONTS:
            if os.path.exists(font_path):
                cls._system_font_path = font_path
                logger.info(f"Found system font: {font_path}")
                return font_path
        
        cls._system_font_path = ""  # 标记为已搜索但未找到
        logger.warning("No CJK system font found")
        return None
    
    @classmethod
    def _get_font(cls, size_pt: float) -> Optional[Any]:
        """获取指定大小的字体对象（带缓存）"""
        if not HAS_PIL:
            return None
        
        # 四舍五入到 0.5pt 以提高缓存效率
        cache_key = round(size_pt * 2) / 2
        
        if cache_key not in cls._font_cache:
            font_path = cls._find_system_font()
            if font_path:
                try:
                    cls._font_cache[cache_key] = ImageFont.truetype(font_path, int(size_pt))
                except Exception as e:
                    logger.warning(f"Failed to load font: {e}")
                    return None
            else:
                # 使用默认字体
                try:
                    cls._font_cache[cache_key] = ImageFont.load_default()
                except:
                    return None
        
        return cls._font_cache.get(cache_key)
    
    @classmethod
    def _measure_text_width(cls, text: str, font_size_pt: float) -> Optional[float]:
        """
        使用实际字体测量文本宽度
        
        Args:
            text: 要测量的文本
            font_size_pt: 字体大小（点）
            
        Returns:
            文本宽度（点），如果测量失败返回 None
        """
        font = cls._get_font(font_size_pt)
        if font is None:
            return None
        
        try:
            bbox = font.getbbox(text)
            width_px = bbox[2] - bbox[0]
            return width_px
        except Exception as e:
            logger.warning(f"Failed to measure text: {e}")
            return None
    
    def __init__(
        self,
        slide_width_inches: float = None,
        slide_height_inches: float = None
    ):
        """
        初始化 PPTX 构建器
        
        Args:
            slide_width_inches: 幻灯片宽度（英寸）
            slide_height_inches: 幻灯片高度（英寸）
        """
        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs: Optional[Presentation] = None
        self.current_slide = None
    
    def create_presentation(self) -> Presentation:
        """创建新的演示文稿"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs
    
    def add_blank_slide(self):
        """添加空白幻灯片"""
        if not self.prs:
            self.create_presentation()
        
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide
    
    def pixels_to_inches(self, pixels: float, dpi: int = None) -> float:
        """像素转英寸"""
        dpi = dpi or self.DEFAULT_DPI
        return pixels / dpi
    
    def calculate_font_size(
        self,
        bbox_width_inches: float,
        bbox_height_inches: float,
        text: str,
        max_font_size: float = None,
        min_font_size: float = None
    ) -> float:
        """
        计算适合边界框的字体大小
        
        Args:
            bbox_width_inches: 边界框宽度（英寸）
            bbox_height_inches: 边界框高度（英寸）
            text: 文本内容
            max_font_size: 最大字体大小
            min_font_size: 最小字体大小
            
        Returns:
            计算出的字体大小（点）
        """
        max_size = max_font_size or self.MAX_FONT_SIZE
        min_size = min_font_size or self.MIN_FONT_SIZE
        
        # 转换为点
        width_pt = bbox_width_inches * 72
        height_pt = bbox_height_inches * 72
        
        if width_pt <= 0 or height_pt <= 0:
            return min_size
        
        # 检查是否可以使用精确测量
        use_precise = HAS_PIL and self._find_system_font()
        
        best_size = min_size
        
        # 二分搜索最大可用字体大小
        for font_size in range(int(max_size), int(min_size) - 1, -1):
            font_size = float(font_size)
            
            # 处理包含换行符的文本
            lines = text.split('\n')
            total_required_lines = 0
            
            for line in lines:
                if not line:
                    total_required_lines += 1
                    continue
                
                # 测量行宽
                if use_precise:
                    line_width_pt = self._measure_text_width(line, font_size)
                    if line_width_pt is None:
                        use_precise = False
                
                if not use_precise:
                    # 回退：基于字符数估算
                    cjk_count = sum(1 for c in line if '\u4e00' <= c <= '\u9fff')
                    non_cjk_count = len(line) - cjk_count
                    line_width_pt = (cjk_count * 1.0 + non_cjk_count * 0.55) * font_size
                
                # 计算此行需要多少行（自动换行）
                if line_width_pt and width_pt > 0:
                    lines_needed = max(1, int(line_width_pt / width_pt) + (1 if line_width_pt % width_pt > 0 else 0))
                else:
                    lines_needed = 1
                total_required_lines += lines_needed
            
            # 计算总高度
            line_height_pt = font_size * 1.2  # 行高系数
            total_height_pt = total_required_lines * line_height_pt
            
            if total_height_pt <= height_pt:
                best_size = font_size
                break
        
        return best_size
    
    def add_background_image(
        self,
        slide,
        image_data: bytes
    ):
        """
        添加背景图片
        
        Args:
            slide: 目标幻灯片
            image_data: 图片字节数据
        """
        try:
            image_stream = BytesIO(image_data)
            slide.shapes.add_picture(
                image_stream,
                Inches(0), Inches(0),
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
        except Exception as e:
            logger.error(f"Failed to add background image: {e}")
    
    def add_title(
        self,
        slide,
        text: str,
        font_size: float = None,
        font_color: str = "#FFFFFF",
        bold: bool = True,
        align: str = "center",
        top_inches: float = 0.5,
        height_inches: float = 1.5
    ):
        """
        添加标题
        
        Args:
            slide: 目标幻灯片
            text: 标题文本
            font_size: 字体大小（点）
            font_color: 字体颜色（十六进制）
            bold: 是否加粗
            align: 对齐方式
            top_inches: 距顶部距离（英寸）
            height_inches: 高度（英寸）
        """
        margin = self.DEFAULT_MARGIN_INCHES
        width_inches = self.slide_width_inches - (margin * 2)
        
        # 自动计算字体大小
        if font_size is None:
            font_size = self.calculate_font_size(
                width_inches, height_inches, text,
                max_font_size=self.DEFAULT_TITLE_FONT_SIZE
            )
        
        textbox = slide.shapes.add_textbox(
            Inches(margin),
            Inches(top_inches),
            Inches(width_inches),
            Inches(height_inches)
        )
        
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        
        # 设置颜色
        r, g, b = self._hex_to_rgb(font_color)
        paragraph.font.color.rgb = RGBColor(r, g, b)
        
        # 设置对齐
        if align == "center":
            paragraph.alignment = PP_ALIGN.CENTER
        elif align == "right":
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        return textbox
    
    def add_content(
        self,
        slide,
        text: str,
        font_size: float = None,
        font_color: str = "#FFFFFF",
        bold: bool = False,
        align: str = "left",
        top_inches: float = 2.5,
        height_inches: float = 4.5
    ):
        """
        添加内容文本
        
        Args:
            slide: 目标幻灯片
            text: 内容文本
            font_size: 字体大小（点）
            font_color: 字体颜色（十六进制）
            bold: 是否加粗
            align: 对齐方式
            top_inches: 距顶部距离（英寸）
            height_inches: 高度（英寸）
        """
        margin = self.DEFAULT_MARGIN_INCHES
        width_inches = self.slide_width_inches - (margin * 2)
        
        # 自动计算字体大小
        if font_size is None:
            font_size = self.calculate_font_size(
                width_inches, height_inches, text,
                max_font_size=self.DEFAULT_BODY_FONT_SIZE,
                min_font_size=12
            )
        
        textbox = slide.shapes.add_textbox(
            Inches(margin),
            Inches(top_inches),
            Inches(width_inches),
            Inches(height_inches)
        )
        
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)
        
        # 处理多行文本
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            
            para.text = line.strip()
            para.font.size = Pt(font_size)
            para.font.bold = bold
            para.space_after = Pt(8)  # 段落间距
            
            # 设置颜色
            r, g, b = self._hex_to_rgb(font_color)
            para.font.color.rgb = RGBColor(r, g, b)
            
            # 设置对齐
            if align == "center":
                para.alignment = PP_ALIGN.CENTER
            elif align == "right":
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
        
        return textbox
    
    def add_subtitle(
        self,
        slide,
        text: str,
        font_size: float = None,
        font_color: str = "#CCCCCC",
        top_inches: float = 2.0
    ):
        """
        添加副标题
        
        Args:
            slide: 目标幻灯片
            text: 副标题文本
            font_size: 字体大小（点）
            font_color: 字体颜色（十六进制）
            top_inches: 距顶部距离（英寸）
        """
        return self.add_content(
            slide, text,
            font_size=font_size or self.DEFAULT_SUBTITLE_FONT_SIZE,
            font_color=font_color,
            align="center",
            top_inches=top_inches,
            height_inches=1.0
        )
    
    def add_notes(self, slide, notes_text: str):
        """添加备注"""
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
    
    def add_illustration(
        self,
        slide,
        image_data: bytes,
        left_inches: float = 7.0,
        top_inches: float = 1.5,
        width_inches: float = 5.5,
        height_inches: float = 5.0
    ):
        """
        添加配图（作为独立图片对象，可编辑/移动）
        
        配图是独立的插图，放置在 PPT 页面的指定位置，
        用户可以自由移动、调整大小。
        
        Args:
            slide: 目标幻灯片
            image_data: 图片字节数据
            left_inches: 左边距（英寸）
            top_inches: 上边距（英寸）
            width_inches: 宽度（英寸）
            height_inches: 高度（英寸）
            
        Returns:
            添加的图片 shape
        """
        try:
            image_stream = BytesIO(image_data)
            picture = slide.shapes.add_picture(
                image_stream,
                Inches(left_inches),
                Inches(top_inches),
                width=Inches(width_inches),
                height=Inches(height_inches)
            )
            logger.info(f"配图已添加: 位置=({left_inches}, {top_inches}), 大小=({width_inches}x{height_inches})")
            return picture
        except Exception as e:
            logger.error(f"添加配图失败: {e}")
            return None
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """十六进制颜色转 RGB"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def save(self, output_path: str):
        """保存演示文稿"""
        if not self.prs:
            raise ValueError("No presentation to save")
        
        output_path_obj = Path(output_path)
        output_dir = output_path_obj.parent
        if str(output_dir) != '.':
            output_dir.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(output_path)
        logger.info(f"Saved presentation to: {output_path}")
    
    def save_to_bytes(self) -> bytes:
        """保存为字节流"""
        if not self.prs:
            raise ValueError("No presentation to save")
        
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    def get_presentation(self) -> Presentation:
        """获取当前演示文稿对象"""
        return self.prs


class SlideBuilder:
    """
    单页幻灯片构建器 - 提供链式 API
    
    用法:
        builder = SlideBuilder(pptx_builder)
        builder.add_background(image_data)
               .add_title("标题")
               .add_content("内容")
               .add_notes("备注")
    """
    
    def __init__(self, pptx_builder: PPTXBuilder, slide=None):
        self.builder = pptx_builder
        self.slide = slide or pptx_builder.add_blank_slide()
    
    def add_background(self, image_data: bytes) -> 'SlideBuilder':
        """添加背景图片"""
        self.builder.add_background_image(self.slide, image_data)
        return self
    
    def add_title(
        self,
        text: str,
        **kwargs
    ) -> 'SlideBuilder':
        """添加标题"""
        self.builder.add_title(self.slide, text, **kwargs)
        return self
    
    def add_content(
        self,
        text: str,
        **kwargs
    ) -> 'SlideBuilder':
        """添加内容"""
        self.builder.add_content(self.slide, text, **kwargs)
        return self
    
    def add_subtitle(
        self,
        text: str,
        **kwargs
    ) -> 'SlideBuilder':
        """添加副标题"""
        self.builder.add_subtitle(self.slide, text, **kwargs)
        return self
    
    def add_notes(self, notes_text: str) -> 'SlideBuilder':
        """添加备注"""
        self.builder.add_notes(self.slide, notes_text)
        return self
    
    def add_illustration(
        self,
        image_data: bytes,
        left_inches: float = 7.0,
        top_inches: float = 1.5,
        width_inches: float = 5.5,
        height_inches: float = 5.0
    ) -> 'SlideBuilder':
        """
        添加配图（独立图片，可编辑/移动）
        
        Args:
            image_data: 图片字节数据
            left_inches: 左边距
            top_inches: 上边距
            width_inches: 宽度
            height_inches: 高度
        """
        self.builder.add_illustration(
            self.slide, image_data,
            left_inches, top_inches,
            width_inches, height_inches
        )
        return self
    
    def get_slide(self):
        """获取幻灯片对象"""
        return self.slide

