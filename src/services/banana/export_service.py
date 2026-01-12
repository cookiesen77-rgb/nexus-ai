"""
Nexus PPT 导出服务
处理 PPTX 和 PDF 导出
"""

import os
import io
import base64
import logging
import uuid
from pathlib import Path
from typing import List, Dict, Optional, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)


class ExportService:
    """导出服务"""
    
    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or "/Users/mac/Desktop/manus/uploads/ppt_exports"
        Path(self.output_dir).mkdir(parents=True, exist_ok=True)
    
    def export_pptx(
        self,
        pages: List[Dict],
        project_name: str = "演示文稿",
        template_style: str = "dark"
    ) -> str:
        """
        导出为 PPTX 文件
        
        Args:
            pages: 页面列表，每个包含 image_url 或 image_base64
            project_name: 项目名称
            template_style: 模板风格
            
        Returns:
            导出文件路径
        """
        # 创建演示文稿
        prs = Presentation()
        
        # 设置幻灯片尺寸为 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 获取空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局
        
        for page in pages:
            slide = prs.slides.add_slide(blank_layout)
            
            # 获取图片
            image_data = page.get("image_base64") or page.get("image_url")
            
            if image_data:
                try:
                    # 如果是 base64
                    if not image_data.startswith("http"):
                        if image_data.startswith("data:"):
                            image_data = image_data.split(",")[1]
                        image_bytes = base64.b64decode(image_data)
                        image_stream = io.BytesIO(image_bytes)
                    else:
                        # 如果是 URL，下载图片
                        import httpx
                        response = httpx.get(image_data)
                        image_stream = io.BytesIO(response.content)
                    
                    # 添加全屏图片
                    slide.shapes.add_picture(
                        image_stream,
                        Inches(0),
                        Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"添加图片失败: {e}")
                    # 添加占位文本
                    self._add_placeholder_text(slide, page.get("title", "幻灯片"), prs)
            else:
                # 没有图片，添加文本内容
                self._add_text_slide(slide, page, prs)
        
        # 保存文件
        filename = f"{project_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        logger.info(f"[ExportService] PPTX 导出成功: {filepath}")
        return filepath
    
    def export_pptx_bytes(
        self,
        pages: List[Dict],
        project_name: str = "演示文稿"
    ) -> bytes:
        """导出为 PPTX 字节流"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_layout = prs.slide_layouts[6]
        
        for page in pages:
            slide = prs.slides.add_slide(blank_layout)
            
            image_data = page.get("image_base64") or page.get("image_url")
            
            if image_data:
                try:
                    if not image_data.startswith("http"):
                        if image_data.startswith("data:"):
                            image_data = image_data.split(",")[1]
                        image_bytes = base64.b64decode(image_data)
                        image_stream = io.BytesIO(image_bytes)
                    else:
                        import httpx
                        response = httpx.get(image_data)
                        image_stream = io.BytesIO(response.content)
                    
                    slide.shapes.add_picture(
                        image_stream,
                        Inches(0),
                        Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"添加图片失败: {e}")
                    self._add_placeholder_text(slide, page.get("title", "幻灯片"), prs)
            else:
                self._add_text_slide(slide, page, prs)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.getvalue()
    
    def export_pdf(
        self,
        pages: List[Dict],
        project_name: str = "演示文稿"
    ) -> str:
        """
        导出为 PDF 文件
        通过将图片合并为 PDF
        
        Args:
            pages: 页面列表
            project_name: 项目名称
            
        Returns:
            导出文件路径
        """
        try:
            import img2pdf
        except ImportError:
            logger.error("img2pdf 未安装，无法导出 PDF")
            raise ImportError("请安装 img2pdf: pip install img2pdf")
        
        images = []
        temp_files = []
        
        for i, page in enumerate(pages):
            image_data = page.get("image_base64") or page.get("image_url")
            
            if image_data:
                try:
                    if not image_data.startswith("http"):
                        if image_data.startswith("data:"):
                            image_data = image_data.split(",")[1]
                        image_bytes = base64.b64decode(image_data)
                    else:
                        import httpx
                        response = httpx.get(image_data)
                        image_bytes = response.content
                    
                    # 保存临时文件
                    temp_path = os.path.join(self.output_dir, f"temp_{uuid.uuid4()}.png")
                    with open(temp_path, "wb") as f:
                        f.write(image_bytes)
                    images.append(temp_path)
                    temp_files.append(temp_path)
                except Exception as e:
                    logger.error(f"处理第 {i+1} 页图片失败: {e}")
        
        if not images:
            raise ValueError("没有可导出的图片")
        
        # 生成 PDF
        filename = f"{project_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.output_dir, filename)
        
        with open(filepath, "wb") as f:
            f.write(img2pdf.convert(images))
        
        # 清理临时文件
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
        
        logger.info(f"[ExportService] PDF 导出成功: {filepath}")
        return filepath
    
    def _add_placeholder_text(self, slide, title: str, prs):
        """添加占位文本"""
        left = Inches(1)
        top = Inches(3)
        width = prs.slide_width - Inches(2)
        height = Inches(1.5)
        
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 添加深色背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(30, 30, 40)
    
    def _add_text_slide(self, slide, page: Dict, prs):
        """添加文本幻灯片"""
        title = page.get("title", "")
        content = page.get("description_content", "")
        
        # 设置背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(30, 30, 40)
        
        # 添加标题
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                prs.slide_width - Inches(1), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 添加内容
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8),
                prs.slide_width - Inches(1), prs.slide_height - Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # 解析内容
            lines = content.split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(220, 220, 220)
                p.space_after = Pt(12)


# 全局导出服务实例
_export_service: Optional[ExportService] = None


def get_export_service() -> ExportService:
    """获取导出服务单例"""
    global _export_service
    if _export_service is None:
        _export_service = ExportService()
    return _export_service

