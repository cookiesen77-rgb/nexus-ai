"""
PPT 核心服务 - 参考 banana-slides 项目的专业实现

核心流程：
1. 大纲生成 - 使用 Gemini-3-pro-preview 生成结构化大纲
2. 页面内容生成 - 使用 Gemini-3-pro-preview 生成文案
3. 配图生成 - 使用 Gemini-2.5-flash-image 生成配图（不是整页图像！）
4. PPTX 导出 - 文案 + 配图组合成可编辑的 PPTX

重要：配图是独立的插图，不是整个 PPT 页面的图像！
"""

import asyncio
import json
import logging
import os
import re
import base64
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Optional, Callable, Any, List, Dict, Union
from PIL import Image

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.models.ppt import (
    Presentation, Slide, SlideLayout, TemplateStyle,
    get_template, get_all_templates, TEMPLATES
)
from src.services.gemini_image import get_gemini_client, GeminiImageClient
from src.services.gemini_chat import get_gemini_chat_client, GeminiChatClient
from src.services.ppt_prompts import (
    get_outline_generation_prompt,
    get_page_description_prompt,
    get_image_generation_prompt,
    get_slide_image_prompt,
    get_illustration_prompt,
    get_language_instruction
)
from src.utils.pptx_builder import PPTXBuilder, SlideBuilder

logger = logging.getLogger(__name__)


class PPTService:
    """
    PPT 生成服务
    
    核心流程：
    1. 大纲生成 - 使用 Gemini-3-pro-preview 生成结构化大纲
    2. 页面内容生成 - 使用 Gemini-3-pro-preview 生成文案
    3. 配图生成 - 使用 Gemini-2.5-flash-image 生成配图（独立插图，不含文字）
    4. PPTX 导出 - 文案 + 配图组合成可编辑的 PPTX
    
    重要：配图是独立的插图，不是整个 PPT 页面的图像！
    """
    
    # 默认配置
    DEFAULT_ASPECT_RATIO = "16:9"
    DEFAULT_RESOLUTION = "4K"
    MAX_CONTENT_LENGTH = 25  # 每条要点最大字数
    
    def __init__(
        self,
        llm_client: Any = None,
        gemini_client: Optional[GeminiImageClient] = None,
        gemini_chat_client: Optional[GeminiChatClient] = None
    ):
        # Gemini Chat 客户端 - 用于大纲和文案生成
        self.gemini_chat = gemini_chat_client or get_gemini_chat_client()
        # Gemini Image 客户端 - 用于配图生成
        self.gemini_image = gemini_client or get_gemini_client()
        # 兼容旧代码
        self.llm_client = llm_client or self.gemini_chat
        self.gemini_client = self.gemini_image
        
        self.presentations: Dict[str, Presentation] = {}
        # 模板图片存储（base64 或文件路径）
        self.template_images: Dict[str, str] = {}
    
    # =========================================================================
    # 模板系统
    # =========================================================================
    def set_template_image(self, presentation_id: str, image_data: Union[str, bytes]) -> bool:
        """
        设置演示文稿的模板参考图片
        
        Args:
            presentation_id: 演示文稿 ID
            image_data: base64 字符串或图片字节
            
        Returns:
            是否设置成功
        """
        try:
            if isinstance(image_data, bytes):
                image_data = base64.b64encode(image_data).decode()
            self.template_images[presentation_id] = image_data
            return True
        except Exception as e:
            logger.error(f"设置模板图片失败: {e}")
            return False
    
    def get_template_image(self, presentation_id: str) -> Optional[Image.Image]:
        """获取模板参考图片的 PIL Image 对象"""
        image_data = self.template_images.get(presentation_id)
        if not image_data:
            return None
        try:
            image_bytes = base64.b64decode(image_data)
            return Image.open(BytesIO(image_bytes))
        except Exception as e:
            logger.error(f"解析模板图片失败: {e}")
            return None
    
    # =========================================================================
    # 大纲生成（使用 Gemini-3-pro-preview）
    # =========================================================================
    async def generate_outline(
        self,
        topic: str,
        page_count: int = 8,
        requirements: str = "",
        language: str = "zh",
        progress_callback: Optional[Callable] = None
    ) -> List[Dict]:
        """
        生成 PPT 大纲（使用 Gemini-3-pro-preview）
        
        Args:
            topic: 演示主题
            page_count: 页数
            requirements: 额外需求
            language: 输出语言
            progress_callback: 进度回调函数
            
        Returns:
            大纲列表，每项包含 title, points, layout
        """
        if progress_callback:
            await progress_callback("generating_outline", 0, page_count, "正在使用 Gemini 生成大纲...")
        
        try:
            # 使用 Gemini Chat 客户端生成大纲
            outline = await self.gemini_chat.generate_outline(
                topic=topic,
                page_count=page_count,
                requirements=requirements,
                language=language
            )
            
            if progress_callback:
                await progress_callback("outline_complete", page_count, page_count, "大纲生成完成")
            
            return outline
            
        except Exception as e:
            logger.error(f"生成大纲失败: {e}")
            return self._generate_default_outline_list(topic, page_count)
    
    def _parse_outline_response(self, text: str, topic: str, page_count: int) -> List[Dict]:
        """解析大纲响应，支持两种格式"""
        try:
            # 清理 markdown 代码块
            text = text.strip()
            if text.startswith("```"):
                text = re.sub(r'^```\w*\n?', '', text)
                text = re.sub(r'\n?```$', '', text)
            
            # 尝试提取 JSON
            json_match = re.search(r'\[[\s\S]*\]', text)
            if json_match:
                outline = json.loads(json_match.group())
                # 展平 part-based 格式
                return self._flatten_outline(outline)
        except json.JSONDecodeError as e:
            logger.warning(f"大纲 JSON 解析失败: {e}")
        
        return self._generate_default_outline_list(topic, page_count)
    
    def _flatten_outline(self, outline: List[Dict]) -> List[Dict]:
        """展平大纲（支持 part-based 格式）"""
        pages = []
        for item in outline:
            if "part" in item and "pages" in item:
                # Part-based 格式
                for page in item["pages"]:
                    page_copy = page.copy()
                    page_copy["part"] = item["part"]
                    pages.append(page_copy)
            else:
                # 简单格式
                pages.append(item)
        return pages
    
    def _generate_default_outline_json(self, topic: str, page_count: int) -> str:
        """生成默认大纲 JSON 字符串"""
        outline = self._generate_default_outline_list(topic, page_count)
        return json.dumps(outline, ensure_ascii=False)
    
    def _generate_default_outline_list(self, topic: str, page_count: int) -> List[Dict]:
        """生成默认大纲列表"""
        outline = [
            {"title": topic, "points": ["演示文稿"], "layout": "title"}
        ]
        
        middle_pages = page_count - 2
        for i in range(middle_pages):
            if i == 0:
                outline.append({
                    "title": "目录",
                    "points": ["第一部分", "第二部分", "第三部分"],
                    "layout": "title_content"
                })
            else:
                outline.append({
                    "title": f"第{i}部分",
                    "points": [f"要点一", f"要点二", f"要点三"],
                    "layout": "title_content"
                })
        
        outline.append({
            "title": "总结",
            "points": [f"{topic}的核心要点", "行动建议", "感谢聆听"],
            "layout": "conclusion"
        })
        
        return outline
    
    # =========================================================================
    # 页面描述生成
    # =========================================================================
    async def generate_page_descriptions(
        self,
        topic: str,
        outline: List[Dict],
        language: str = "zh",
        progress_callback: Optional[Callable] = None
    ) -> List[str]:
        """
        为所有页面生成详细描述
        
        Args:
            topic: 主题
            outline: 大纲
            language: 语言
            progress_callback: 进度回调
            
        Returns:
            页面描述列表
        """
        descriptions = []
        previous_context = ""
        
        for i, page_outline in enumerate(outline):
            page_index = i + 1
            
            if progress_callback:
                await progress_callback(
                    "generating_description",
                    i, len(outline),
                    f"正在生成第 {page_index}/{len(outline)} 页描述..."
                )
            
            prompt = get_page_description_prompt(
                topic=topic,
                outline=outline,
                page_outline=page_outline,
                page_index=page_index,
                previous_context=previous_context,
                language=language
            )
            
            try:
                if self.llm_client:
                    response = await self.llm_client.complete([
                        {"role": "user", "content": prompt}
                    ])
                    description = response.get("content", "")
                else:
                    # 默认描述
                    description = self._generate_default_description(page_outline, page_index)
                
                # 优化描述（确保简洁）
                description = self._optimize_description(description)
                descriptions.append(description)
                
                # 更新上下文
                previous_context += f"\n第{page_index}页: {page_outline.get('title', '')}"
                
            except Exception as e:
                logger.error(f"生成第 {page_index} 页描述失败: {e}")
                descriptions.append(self._generate_default_description(page_outline, page_index))
        
        return descriptions
    
    def _generate_default_description(self, page_outline: Dict, page_index: int) -> str:
        """生成默认页面描述"""
        title = page_outline.get("title", "未命名")
        points = page_outline.get("points", [])
        
        description = f"页面标题：{title}\n\n页面文字：\n"
        for point in points[:5]:  # 最多 5 个要点
            description += f"- {point[:self.MAX_CONTENT_LENGTH]}\n"
        
        return description
    
    def _optimize_description(self, description: str) -> str:
        """优化描述，确保简洁"""
        lines = description.split('\n')
        optimized_lines = []
        
        for line in lines:
            line = line.strip()
            if line.startswith('-') or line.startswith('•'):
                # 截断过长的要点
                content = line[1:].strip()
                if len(content) > self.MAX_CONTENT_LENGTH:
                    content = content[:self.MAX_CONTENT_LENGTH - 3] + "..."
                optimized_lines.append(f"- {content}")
            else:
                optimized_lines.append(line)
        
        return '\n'.join(optimized_lines)
    
    # =========================================================================
    # 配图生成（使用 Gemini-2.5-flash-image）
    # 重要：这是生成配图，不是整个 PPT 页面的图像！
    # =========================================================================
    async def generate_slide_illustrations(
        self,
        presentation: Presentation,
        slides_content: List[Dict] = None,
        progress_callback: Optional[Callable] = None
    ) -> Presentation:
        """
        为演示文稿的幻灯片生成配图（独立插图，不含文字）
        
        配图会被放置在 PPT 页面的插图区域，与文字内容并排显示。
        
        Args:
            presentation: 演示文稿对象
            slides_content: 幻灯片内容列表（包含 needs_illustration, illustration_theme）
            progress_callback: 进度回调
            
        Returns:
            更新后的演示文稿
        """
        total = len(presentation.slides)
        template_config = get_template(presentation.template)
        
        # 确定风格
        style = "professional"
        if template_config:
            style_map = {
                "modern": "professional",
                "minimal": "minimal",
                "creative": "creative",
                "nature": "nature",
                "dark": "tech"
            }
            style = style_map.get(presentation.template, "professional")
        
        for i, slide in enumerate(presentation.slides):
            # 检查是否需要配图
            slide_info = slides_content[i] if slides_content and i < len(slides_content) else {}
            needs_illustration = slide_info.get("needs_illustration", True)
            
            # 封面页和结尾页通常不需要配图
            is_cover = (i == 0 or slide.layout == SlideLayout.TITLE.value)
            is_conclusion = (i == total - 1 or slide.layout == SlideLayout.CONCLUSION.value)
            
            if not needs_illustration or is_cover or is_conclusion:
                logger.info(f"幻灯片 {i+1} 跳过配图生成")
                slide.image_base64 = ""  # 清空，表示不需要配图
                continue
            
            if progress_callback:
                await progress_callback(
                    "generating_illustration",
                    i, total,
                    f"正在为第 {i+1}/{total} 页生成配图..."
                )
            
            # 获取配图主题
            illustration_theme = slide_info.get("illustration_theme", slide.title)
            
            # 调用 Gemini Image 生成配图
            result = await self.gemini_image.generate_illustration(
                topic=presentation.topic,
                slide_title=slide.title,
                slide_content=slide.content,
                illustration_theme=illustration_theme,
                style=style,
                aspect_ratio=self.DEFAULT_ASPECT_RATIO
            )
            
            if result.get("success"):
                slide.image_base64 = result["image_base64"]
                slide.image_prompt = f"配图主题: {illustration_theme}"
                logger.info(f"幻灯片 {i+1} 配图生成成功")
            else:
                logger.warning(f"幻灯片 {i+1} 配图生成失败: {result.get('error')}")
                slide.image_base64 = ""
            
            # 短暂延迟，避免 API 限流
            await asyncio.sleep(0.5)
        
        if progress_callback:
            await progress_callback("illustrations_complete", total, total, "配图生成完成")
        
        presentation.update_timestamp()
        return presentation
    
    # 保留旧方法名以兼容
    async def generate_slide_images(
        self,
        presentation: Presentation,
        progress_callback: Optional[Callable] = None
    ) -> Presentation:
        """兼容旧代码的别名"""
        return await self.generate_slide_illustrations(
            presentation, 
            progress_callback=progress_callback
        )
    
    def _generate_outline_text(self, slides: List[Slide]) -> str:
        """生成大纲文本"""
        lines = []
        for i, slide in enumerate(slides, 1):
            lines.append(f"{i}. {slide.title}")
        return "\n".join(lines)
    
    # =========================================================================
    # 完整创建流程
    # =========================================================================
    async def create_presentation(
        self,
        topic: str,
        page_count: int = 8,
        template: str = "modern",
        requirements: str = "",
        template_image: Optional[bytes] = None,
        progress_callback: Optional[Callable] = None
    ) -> Presentation:
        """
        创建完整的演示文稿
        
        流程：
        1. 使用 Gemini-3-pro-preview 生成大纲
        2. 使用 Gemini-3-pro-preview 生成页面文案
        3. 使用 Gemini-2.5-flash-image 生成配图（独立插图，不是整页图像）
        
        Args:
            topic: 主题
            page_count: 页数
            template: 模板风格
            requirements: 额外需求
            template_image: 模板参考图片（可选，用于风格参考）
            progress_callback: 进度回调
            
        Returns:
            创建的演示文稿
        """
        # 创建演示文稿对象
        presentation = Presentation(
            title=topic,
            topic=topic,
            template=template
        )
        
        # 设置模板图片（如果提供）
        if template_image:
            self.set_template_image(presentation.id, template_image)
        
        # 1. 使用 Gemini-3-pro-preview 生成大纲
        if progress_callback:
            await progress_callback("generating_outline", 0, page_count, "正在使用 Gemini 生成大纲...")
        
        outline = await self.generate_outline(
            topic, page_count, requirements, "zh", None
        )
        
        if progress_callback:
            await progress_callback("outline_complete", 1, page_count, f"大纲生成完成，共 {len(outline)} 页")
        
        # 2. 使用 Gemini-3-pro-preview 生成页面文案
        slides_content = []
        for i, item in enumerate(outline):
            if progress_callback:
                await progress_callback(
                    "generating_content",
                    i + 1, len(outline),
                    f"正在生成第 {i+1}/{len(outline)} 页文案..."
                )
            
            # 使用 Gemini Chat 生成页面内容
            page_content = await self.gemini_chat.generate_page_content(
                topic=topic,
                page_title=item.get("title", ""),
                page_points=item.get("points", []),
                page_index=i + 1,
                total_pages=len(outline),
                language="zh"
            )
            slides_content.append(page_content)
            
            # 创建幻灯片
            slide = Slide(
                order=i,
                layout=self._map_layout(item.get("layout", "title_content")),
                title=page_content.get("title", item.get("title", "")),
                content=page_content.get("content", "")
            )
            presentation.slides.append(slide)
        
        if progress_callback:
            await progress_callback("content_complete", len(outline), len(outline), "文案生成完成")
        
        # 3. 使用 Gemini-2.5-flash-image 生成配图
        presentation = await self.generate_slide_illustrations(
            presentation, 
            slides_content,
            progress_callback
        )
        
        # 保存
        self.presentations[presentation.id] = presentation
        
        if progress_callback:
            await progress_callback("complete", page_count, page_count, "PPT 创建完成！")
        
        return presentation
    
    def _extract_content_from_description(self, description: str) -> str:
        """从描述中提取页面文字部分"""
        if "页面文字：" in description:
            parts = description.split("页面文字：")
            if len(parts) > 1:
                return parts[1].strip()
        return description
    
    def _map_layout(self, layout: str) -> str:
        """映射布局名称"""
        layout_map = {
            "title": SlideLayout.TITLE.value,
            "title_content": SlideLayout.TITLE_CONTENT.value,
            "title_image": SlideLayout.TITLE_IMAGE.value,
            "section": SlideLayout.SECTION.value,
            "conclusion": SlideLayout.CONCLUSION.value,
        }
        return layout_map.get(layout, SlideLayout.TITLE_CONTENT.value)
    
    # =========================================================================
    # 其他方法
    # =========================================================================
    async def regenerate_slide_image(
        self,
        presentation_id: str,
        slide_index: int,
        custom_prompt: Optional[str] = None
    ) -> Optional[Slide]:
        """重新生成单张幻灯片的配图"""
        presentation = self.presentations.get(presentation_id)
        if not presentation or slide_index >= len(presentation.slides):
            return None
        
        slide = presentation.slides[slide_index]
        ref_image = self.get_template_image(presentation_id)
        
        if custom_prompt:
            prompt = custom_prompt
        else:
            template_config = get_template(presentation.template)
            outline_text = self._generate_outline_text(presentation.slides)
            prompt = get_image_generation_prompt(
                page_desc=f"页面标题：{slide.title}\n\n页面文字：\n{slide.content}",
                outline_text=outline_text,
                current_section=slide.title,
                page_index=slide_index + 1,
                has_template=ref_image is not None,
                template_style=template_config.description if template_config else "",
            )
        
        result = await self.gemini_client.generate_image(
            prompt=prompt,
            ref_images=[ref_image] if ref_image else None,
            aspect_ratio=self.DEFAULT_ASPECT_RATIO
        )
        
        if result.get("success"):
            slide.image_base64 = result["image_base64"]
            slide.image_prompt = prompt[:500]
            presentation.update_timestamp()
        
        return slide
    
    def update_slide(
        self,
        presentation_id: str,
        slide_index: int,
        updates: Dict
    ) -> Optional[Slide]:
        """更新幻灯片内容"""
        presentation = self.presentations.get(presentation_id)
        if not presentation or slide_index >= len(presentation.slides):
            return None
        
        slide = presentation.slides[slide_index]
        
        if "title" in updates:
            slide.title = updates["title"]
        if "content" in updates:
            slide.content = updates["content"]
        if "layout" in updates:
            slide.layout = updates["layout"]
        if "notes" in updates:
            slide.notes = updates["notes"]
        
        presentation.update_timestamp()
        return slide
    
    def get_presentation(self, presentation_id: str) -> Optional[Presentation]:
        """获取演示文稿"""
        return self.presentations.get(presentation_id)
    
    def delete_presentation(self, presentation_id: str) -> bool:
        """删除演示文稿"""
        if presentation_id in self.presentations:
            del self.presentations[presentation_id]
            if presentation_id in self.template_images:
                del self.template_images[presentation_id]
            return True
        return False
    
    # =========================================================================
    # PPTX 导出 - 配图作为插图（可编辑）
    # =========================================================================
    def export_pptx(
        self,
        presentation_id: str,
        output_path: Optional[str] = None
    ) -> Optional[str]:
        """
        导出为 PPTX 文件
        
        所有元素都是独立的、可编辑的：
        - 文本框可以移动、编辑
        - 配图作为独立图片对象，可以移动、调整大小
        - 用户可以自由调整布局
        """
        presentation = self.presentations.get(presentation_id)
        if not presentation:
            return None
        
        template_config = get_template(presentation.template)
        
        # 使用专业构建器
        builder = PPTXBuilder()
        builder.create_presentation()
        
        # 获取模板颜色
        text_color = "#333333"  # 默认深色文字（适合浅色背景）
        bg_color = "#FFFFFF"    # 默认白色背景
        if template_config:
            text_color = template_config.colors.get("text", "#333333")
            bg_color = template_config.colors.get("background", "#FFFFFF")
        
        for i, slide_data in enumerate(presentation.slides):
            slide_builder = SlideBuilder(builder)
            
            # 根据布局类型添加内容
            is_title_slide = (
                i == 0 or 
                slide_data.layout == SlideLayout.TITLE.value or 
                slide_data.layout == "title"
            )
            
            is_conclusion_slide = (
                i == len(presentation.slides) - 1 or
                slide_data.layout == SlideLayout.CONCLUSION.value or
                slide_data.layout == "conclusion"
            )
            
            # 判断是否有配图
            has_illustration = bool(slide_data.image_base64)
            
            if is_title_slide:
                # 封面页 - 标题居中，较大字体
                if slide_data.title:
                    slide_builder.add_title(
                        slide_data.title,
                        font_size=56,
                        font_color=text_color,
                        align="center",
                        top_inches=2.5,
                        height_inches=2.0
                    )
                
                # 副标题或内容作为副标题
                if slide_data.content:
                    first_line = slide_data.content.split('\n')[0].strip('- ')
                    slide_builder.add_subtitle(
                        first_line,
                        font_color=self._lighten_color(text_color),
                        top_inches=4.8
                    )
            
            elif is_conclusion_slide:
                # 结尾页
                if slide_data.title:
                    slide_builder.add_title(
                        slide_data.title,
                        font_size=48,
                        font_color=text_color,
                        align="center",
                        top_inches=2.0,
                        height_inches=1.5
                    )
                
                if slide_data.content:
                    slide_builder.add_content(
                        slide_data.content,
                        font_color=text_color,
                        align="center",
                        top_inches=4.0,
                        height_inches=2.5
                    )
            
            else:
                # 普通内容页 - 左侧文字，右侧配图
                if has_illustration:
                    # 有配图：左侧放文字，右侧放配图
                    # 标题
                    if slide_data.title:
                        slide_builder.add_title(
                            slide_data.title,
                            font_size=40,
                            font_color=text_color,
                            align="left",
                            top_inches=0.5,
                            height_inches=1.0
                        )
                    
                    # 内容 - 左侧区域 (约 55% 宽度)
                    if slide_data.content:
                        formatted_content = self._format_bullet_points(slide_data.content)
                        # 使用自定义宽度的内容区域
                        self._add_content_with_width(
                            slide_builder.slide, builder,
                            formatted_content,
                            font_color=text_color,
                            left_inches=0.5,
                            top_inches=1.8,
                            width_inches=6.5,  # 左侧约一半宽度
                            height_inches=5.0
                        )
                    
                    # 配图 - 右侧区域
                    try:
                        image_bytes = base64.b64decode(slide_data.image_base64)
                        slide_builder.add_illustration(
                            image_bytes,
                            left_inches=7.5,   # 右侧
                            top_inches=1.5,
                            width_inches=5.3,
                            height_inches=5.0
                        )
                    except Exception as e:
                        logger.warning(f"添加配图失败: {e}")
                else:
                    # 无配图：全宽文字
                    if slide_data.title:
                        slide_builder.add_title(
                            slide_data.title,
                            font_size=44,
                            font_color=text_color,
                            align="left",
                            top_inches=0.5,
                            height_inches=1.2
                        )
                    
                    if slide_data.content:
                        formatted_content = self._format_bullet_points(slide_data.content)
                        slide_builder.add_content(
                            formatted_content,
                            font_color=text_color,
                            align="left",
                            top_inches=2.0,
                            height_inches=5.0
                        )
            
            # 添加备注
            if slide_data.notes:
                slide_builder.add_notes(slide_data.notes)
        
        # 保存
        if not output_path:
            output_dir = Path(os.environ.get("WORKSPACE_PATH", ".")) / "exports"
            output_dir.mkdir(exist_ok=True)
            output_path = str(output_dir / f"{presentation.id}.pptx")
        
        builder.save(output_path)
        logger.info(f"PPT 已导出到: {output_path}")
        
        return output_path
    
    def _add_content_with_width(
        self,
        slide,
        builder: PPTXBuilder,
        text: str,
        font_color: str,
        left_inches: float,
        top_inches: float,
        width_inches: float,
        height_inches: float
    ):
        """添加自定义宽度的内容文本框"""
        textbox = slide.shapes.add_textbox(
            Inches(left_inches),
            Inches(top_inches),
            Inches(width_inches),
            Inches(height_inches)
        )
        
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 处理多行文本
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            
            para.text = line.strip()
            para.font.size = Pt(24)
            para.space_after = Pt(8)
            
            # 设置颜色
            r, g, b = self._hex_to_rgb(font_color)
            para.font.color.rgb = RGBColor(r, g, b)
    
    def _format_bullet_points(self, content: str) -> str:
        """格式化要点，确保一致的格式"""
        lines = content.split('\n')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 移除现有的项目符号并添加统一格式
            if line.startswith(('- ', '• ', '· ', '* ')):
                line = line[2:]
            elif line.startswith(('-', '•', '·', '*')):
                line = line[1:]
            
            formatted_lines.append(f"• {line.strip()}")
        
        return '\n'.join(formatted_lines)
    
    def _lighten_color(self, hex_color: str) -> str:
        """使颜色变浅"""
        hex_color = hex_color.lstrip('#')
        r, g, b = [int(hex_color[i:i+2], 16) for i in (0, 2, 4)]
        
        # 混合白色使颜色变浅
        factor = 0.3
        r = int(r + (255 - r) * factor)
        g = int(g + (255 - g) * factor)
        b = int(b + (255 - b) * factor)
        
        return f"#{r:02x}{g:02x}{b:02x}"
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """将十六进制颜色转为 RGB"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def export_pptx_bytes(self, presentation_id: str) -> Optional[bytes]:
        """导出为 PPTX 字节流"""
        presentation = self.presentations.get(presentation_id)
        if not presentation:
            return None
        
        template_config = get_template(presentation.template)
        
        # 使用专业构建器
        builder = PPTXBuilder()
        builder.create_presentation()
        
        text_color = "#FFFFFF"
        if template_config:
            text_color = template_config.colors.get("text", "#FFFFFF")
        
        for i, slide_data in enumerate(presentation.slides):
            slide_builder = SlideBuilder(builder)
            
            if slide_data.image_base64:
                try:
                    image_bytes = base64.b64decode(slide_data.image_base64)
                    slide_builder.add_background(image_bytes)
                except Exception as e:
                    logger.warning(f"添加背景图片失败: {e}")
            
            is_title_slide = i == 0 or slide_data.layout in ["title", SlideLayout.TITLE.value]
            
            if is_title_slide:
                if slide_data.title:
                    slide_builder.add_title(
                        slide_data.title, font_size=56,
                        font_color=text_color, align="center",
                        top_inches=2.5, height_inches=2.0
                    )
            else:
                if slide_data.title:
                    slide_builder.add_title(
                        slide_data.title, font_size=44,
                        font_color=text_color, align="left",
                        top_inches=0.5, height_inches=1.2
                    )
                if slide_data.content:
                    formatted_content = self._format_bullet_points(slide_data.content)
                    slide_builder.add_content(
                        formatted_content, font_color=text_color,
                        align="left", top_inches=2.0, height_inches=5.0
                    )
            
            if slide_data.notes:
                slide_builder.add_notes(slide_data.notes)
        
        return builder.save_to_bytes()


# 全局服务实例
_ppt_service: Optional[PPTService] = None


def get_ppt_service() -> PPTService:
    """获取 PPT 服务单例"""
    global _ppt_service
    if _ppt_service is None:
        _ppt_service = PPTService()
    return _ppt_service
